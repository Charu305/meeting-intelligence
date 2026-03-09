from pptx import Presentation

def read_slides(ppt_file):
    prs = Presentation(ppt_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text.append(shape.text)
    return "\n".join(text)